#!/usr/bin/env python3
import argparse
import os
from pathlib import Path
from typing import List, Tuple

from PIL import Image
from natsort import natsorted

# PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# PDF
from reportlab.lib.pagesizes import letter, A4, landscape, portrait
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

SUPPORTED_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}

def collect_images(folder: Path, sort: str) -> List[Path]:
    imgs = [p for p in folder.iterdir() if p.suffix.lower() in SUPPORTED_EXTS]
    if sort == "numeric":
        return natsorted(imgs, key=lambda p: p.name)
    else:
        return sorted(imgs, key=lambda p: p.name.lower())

def best_pagesize_for_image(img_size_px: Tuple[int, int], prefer="letter"):
    """Choose page size orientation per image."""
    w, h = img_size_px
    if prefer == "a4":
        base = A4
    else:
        base = letter
    return landscape(base) if w >= h else portrait(base)

def add_slide_with_image(prs: Presentation, img_path: Path, margin_pts: float, max_dpi: int):
    # Use blank slide layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Determine slide size (in EMUs) and convert to points
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    # EMU to inches: 914400 emu per inch; inches to points: *72
    def emu_to_points(e): return (e / 914400.0) * 72.0

    slide_w_pt = emu_to_points(slide_width_emu)
    slide_h_pt = emu_to_points(slide_height_emu)

    # Load image to know its size
    with Image.open(img_path) as im:
        im = im.convert("RGB")
        img_w_px, img_h_px = im.size

        # compute target box within slide, honoring margins
        box_w_pt = max(1.0, slide_w_pt - 2 * margin_pts)
        box_h_pt = max(1.0, slide_h_pt - 2 * margin_pts)

        # image DPI for scaling cap
        dpi_x = dpi_y = 72
        if "dpi" in im.info:
            try:
                dx, dy = im.info["dpi"]
                dpi_x, dpi_y = int(dx), int(dy)
            except Exception:
                pass

        # limit effective DPI when shrinking a huge image to avoid bloating
        # compute scale-to-fit in points, then ensure we don't exceed max_dpi
        img_aspect = img_w_px / img_h_px
        box_aspect = box_w_pt / box_h_pt
        if img_aspect >= box_aspect:
            # width-bound
            render_w_pt = box_w_pt
            render_h_pt = box_w_pt / img_aspect
        else:
            render_h_pt = box_h_pt
            render_w_pt = box_h_pt * img_aspect

        # left/top offsets to center
        left_pt = (slide_w_pt - render_w_pt) / 2.0
        top_pt = (slide_h_pt - render_h_pt) / 2.0

        # Save a temp downscaled image if necessary to keep DPI reasonable
        # effective DPI = pixels / inches
        target_w_in = render_w_pt / 72.0
        target_h_in = render_h_pt / 72.0
        eff_dpi_x = img_w_px / max(1e-6, target_w_in)
        eff_dpi_y = img_h_px / max(1e-6, target_h_in)
        need_downscale = (eff_dpi_x > max_dpi) or (eff_dpi_y > max_dpi)

        tmp_path = None
        if need_downscale:
            scale = min(max_dpi / eff_dpi_x, max_dpi / eff_dpi_y)
            new_w = max(1, int(img_w_px * scale))
            new_h = max(1, int(img_h_px * scale))
            im = im.resize((new_w, new_h), Image.LANCZOS)
            tmp_path = img_path.with_suffix(".pptx_tmp.jpg")
            im.save(tmp_path, quality=92)

        # Add picture (python-pptx uses EMUs)
        from pptx.util import Pt as _Pt
        from pptx.util import Inches as _Inches

        def pt_to_emu(pt): return int((pt / 72.0) * 914400.0)

        slide.shapes.add_picture(
            str(tmp_path if tmp_path else img_path),
            left=pt_to_emu(left_pt),
            top=pt_to_emu(top_pt),
            width=pt_to_emu(render_w_pt),
            height=pt_to_emu(render_h_pt),
        )

        # Optional: add a tiny filename label
        txbox = slide.shapes.add_textbox(pt_to_emu(8), pt_to_emu(8), pt_to_emu(slide_w_pt-16), pt_to_emu(20))
        tf = txbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = img_path.name
        p.font.size = Pt(8)
        p.alignment = PP_ALIGN.LEFT

        # cleanup temp
        if tmp_path and tmp_path.exists():
            try:
                tmp_path.unlink()
            except Exception:
                pass

def build_pptx(images: List[Path], out_path: Path, margin_pts: float, max_dpi: int):
    # Choose slide size based on the first image orientation
    prs = Presentation()
    # Default slide size is 10" x 7.5" (9144000 x 6858000 EMU). We'll keep that
    # but orient per image page-by-page by centering; PowerPoint slide size is global.

    for img in images:
        add_slide_with_image(prs, img, margin_pts, max_dpi)

    prs.save(str(out_path))

def draw_image_on_pdf_page(c: canvas.Canvas, img: Image.Image, page_w: float, page_h: float, margin: float):
    img_w, img_h = img.size
    box_w = max(1.0, page_w - 2 * margin)
    box_h = max(1.0, page_h - 2 * margin)
    img_aspect = img_w / img_h
    box_aspect = box_w / box_h
    if img_aspect >= box_aspect:
        render_w = box_w
        render_h = box_w / img_aspect
    else:
        render_h = box_h
        render_w = box_h * img_aspect

    x = (page_w - render_w) / 2.0
    y = (page_h - render_h) / 2.0

    c.drawImage(ImageReader(img), x, y, width=render_w, height=render_h, preserveAspectRatio=True, anchor='c')

def build_pdf(images: List[Path], out_path: Path, margin_pts: float, prefer_size: str):
    # Create canvas with a dummy page; we'll set size per page based on image orientation
    c = canvas.Canvas(str(out_path), pagesize=letter)

    for idx, p in enumerate(images, 1):
        with Image.open(p) as im:
            im = im.convert("RGB")
            page_size = best_pagesize_for_image(im.size, prefer="a4" if prefer_size.lower()=="a4" else "letter")
            c.setPageSize(page_size)
            page_w, page_h = page_size
            draw_image_on_pdf_page(c, im, page_w, page_h, margin_pts)

            # footer filename
            c.setFont("Helvetica", 8)
            c.drawString(12, 12, f"{idx}/{len(images)} — {p.name}")

            c.showPage()

    c.save()

def main():
    ap = argparse.ArgumentParser(description="Build a PPTX (one slide per image) and a multi-page PDF from a folder of images.")
    ap.add_argument("folder", type=str, help="Folder containing images")
    ap.add_argument("--pptx", type=str, default=None, help="Output PPTX path (default: images.pptx in folder)")
    ap.add_argument("--pdf", type=str, default=None, help="Output PDF path (default: images.pdf in folder)")
    ap.add_argument("--margin", type=float, default=36.0, help="Margin in points (1 in = 72 pt). Default: 36 (0.5 in)")
    ap.add_argument("--max-dpi", type=int, default=300, help="Max effective DPI when downscaling very large images (PPTX). Default: 300")
    ap.add_argument("--sort", choices=["numeric", "alpha"], default="numeric", help="Sort mode for filenames. Default: numeric")
    ap.add_argument("--pdf-size", choices=["letter", "a4"], default="letter", help="Preferred PDF page size family. Default: letter")
    args = ap.parse_args()

    folder = Path(args.folder).expanduser().resolve()
    if not folder.is_dir():
        raise SystemExit(f"Folder not found: {folder}")

    images = collect_images(folder, sort=args.sort)
    if not images:
        raise SystemExit(f"No images found in {folder} (supported: {', '.join(sorted(SUPPORTED_EXTS))})")

    pptx_path = Path(args.pptx) if args.pptx else folder / "images.pptx"
    pdf_path  = Path(args.pdf)  if args.pdf  else folder / "images.pdf"

    print(f"Found {len(images)} images.")
    print(f"Writing PPTX: {pptx_path}")
    build_pptx(images, pptx_path, margin_pts=args.margin, max_dpi=args.max_dpi)

    print(f"Writing PDF:  {pdf_path}")
    build_pdf(images, pdf_path, margin_pts=args.margin, prefer_size=args.pdf_size)

    print("Done.")

if __name__ == "__main__":
    main()
